from django.shortcuts import render, redirect, get_object_or_404
from django.conf import settings
from django.contrib.auth.decorators import login_required
from django.contrib.auth.models import Group
from django.db.models import Q
from django.http import JsonResponse, HttpResponseForbidden, FileResponse, HttpResponse, HttpResponseBadRequest
from django.contrib import messages
from django.contrib.auth import update_session_auth_hash
from django.shortcuts import redirect, render, get_object_or_404
from django.urls import reverse
from django.utils import timezone
from django.utils.dateparse import parse_date
from django.views.decorators.http import require_GET, require_POST
from django.template.loader import render_to_string
from .models import AccidentReport, Business, TrainingLocation, CourseType, Personnel, Booking, BookingDay, DelegateRegister, FeedbackResponse
from .forms import (
    AccidentReportForm,
    BookingForm,
    PersonnelAdminForm,
    PersonnelProfileForm,
    DelegateRegisterForm,
    FeedbackForm,
)
from .forms_profile import UserProfileForm, PersonnelProfileForm
from .signal_control import disable, enable

from datetime import timedelta, datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE as SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4, portrait
from reportlab.pdfgen import canvas
from reportlab.lib.units import mm
import uuid, io

# NEW: a shape enum that works across python-pptx versions
try:
    # newer python-pptx
    from pptx.enum.shapes import MSO_SHAPE as SHAPE
except Exception:
    # older python-pptx
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE

# --- helpers --------------------------------------------------

def is_instructor_user(user):
    return user.is_authenticated and (
        user.groups.filter(name__iexact="instructor").exists()
        or hasattr(user, "instructor")
    )

def get_instructor_for_user(user):
    """
    Resolve the Instructor profile for this user.
    If you allow multiple mappings, adjust here; otherwise one-to-one.
    """
    try:
        if hasattr(user, "instructor") and user.instructor:
            return user.instructor
    except Instructor.DoesNotExist:
        pass
    # fallback: name/email match, if you’ve used that pattern in your data
    return Personnel.objects.filter(user=user).first()

from datetime import datetime

def _parse_yyyy_mm_dd(s: str):
    """
    Convert 'YYYY-MM-DD' → date object.
    Returns None if invalid.
    """
    if not s:
        return None
    try:
        return datetime.strptime(s, "%Y-%m-%d").date()
    except ValueError:
        return None

# --- home redirect --------------------------------------------

def home(request):
    if request.user.is_authenticated:
        if request.user.is_superuser or request.user.groups.filter(name__iexact="admin").exists():
            return redirect("app_admin_dashboard")
        if is_instructor_user(request.user):
            return redirect("instructor_bookings")
    return render(request, "admin/admin_dashboard.html")

def ensure_groups():
    Group.objects.get_or_create(name='admin')
    Group.objects.get_or_create(name='instructor')

def roles_for(user):
    ensure_groups()
    r=[]
    if user.is_authenticated:
        if user.groups.filter(name='admin').exists(): r.append('admin')
        if user.groups.filter(name='instructor').exists(): r.append('instructor')
        if user.is_superuser: r.append('superuser')
    return r

@require_GET
def public_delegate_instructors_api(request):
    """
    JSON:  /register/instructors?ct=EFAAW&date=2025-10-21
    Returns [{id, name}, ...] for the instructor select.
    """
    code = (request.GET.get("ct") or request.GET.get("code") or "").strip()
    date_str = (request.GET.get("date") or request.GET.get("d") or "").strip()
    day_date = _parse_yyyy_mm_dd(date_str) or timezone.localdate()

    course = CourseType.objects.filter(code__iexact=code).first()
    if not course:
        return JsonResponse({"instructors": []})

    qs = (
        Personnel.objects.filter(
            bookings__course_type=course,
            bookings__days__date=day_date
        )
        .distinct()
        .order_by("name")
    )

    data = [{"id": i.id, "name": i.name} for i in qs]
    return JsonResponse({"instructors": data})


def public_delegate_register(request):
    """
    Public delegate register page.
    - Optional ?ct=<COURSE_CODE> to preselect a course
    - Optional ?date=YYYY-MM-DD to preselect the date
    - Instructor list is restricted to instructors who have a booking
      for the selected course and date.
    """
    # 1) Resolve course from ?ct=<code> or posted course_id
    course = None
    ct_code = (request.GET.get("ct") or "").strip()
    if ct_code:
        course = CourseType.objects.filter(code__iexact=ct_code).first()

    # For the dropdown when no QR is used
    course_types = CourseType.objects.order_by("name")

    # 2) Initial date (from querystring if present, else today)
    initial = {}
    if request.GET.get("date"):
        initial["date"] = request.GET["date"]              # 'YYYY-MM-DD'
    else:
        initial["date"] = timezone.localdate()             # today (date obj)

    form = DelegateRegisterForm(request.POST or None, initial=initial)

    # 3) Build instructors list (depends on course + date)
    instructors = []
    bound_date = form.data.get("date") if form.is_bound else initial.get("date")
    if course and bound_date:
        instructors = (
            Personnel.objects
            .filter(
                bookings__course_type=course,
                bookings__days__date=bound_date
            )
            .distinct()
            .order_by("name")
        )

    if request.method == "POST" and form.is_valid():
        delegate = form.save(commit=False)

        inst = delegate.instructor
        bd = None
        if course and inst and form.cleaned_data.get("date"):
            bd = (
                BookingDay.objects
                .filter(
                    booking__course_type=course,
                    booking__instructor=inst,
                    date=form.cleaned_data["date"],
                )
                .first()
            )
        delegate.booking_day = bd
        delegate.save()

        messages.success(request, "Thank you — you’re on the register.")
        return redirect(request.get_full_path())

    return render(
        request,
        "public/delegate_register.html",
        {
            "form": form,
            "course": course,
            "course_types": course_types,
            "instructors": instructors,
        },
    )

@login_required
def switch_role(request, role):
    if role in roles_for(request.user):
        request.session['current_role']=role
        messages.success(request, f"Switched to {role}")
    else:
        messages.error(request, "You don't have that role.")
    return redirect('home')

@login_required
def home(request):
    if 'current_role' not in request.session:
        r = roles_for(request.user)
        if r: request.session['current_role']=r[0]
    return render(request,'admin/admin_dashboard.html',{})

# Admin app
@login_required
def app_admin_dashboard(request):
    if not (request.user.is_superuser or request.user.groups.filter(name='admin').exists()):
        return HttpResponseForbidden()
    bookings = Booking.objects.select_related('business','training_location','course_type','instructor').order_by('-course_date')[:50]
    return render(request,'admin_app/dashboard.html',{'bookings':bookings})

@login_required
def app_admin_booking_new(request):
    if not (request.user.is_superuser or request.user.groups.filter(name='admin').exists()):
        return HttpResponseForbidden()
    form = BookingForm(request.POST or None)
    if request.method=='POST' and form.is_valid():
        b = form.save()
        n = int(float(b.course_duration_days or b.course_type.duration_days) + 0.999)
        from datetime import timedelta
        for i in range(1,n+1):
            BookingDay.objects.create(booking=b, day_no=i, day_date=b.course_date+timedelta(days=i-1), start_time=b.start_time)
        messages.success(request,"Booking created")
        return redirect('app_admin_booking_detail', pk=b.id)
    return render(request,'admin_app/booking_form.html',{'form':form})

@login_required
def app_admin_booking_detail(request, pk):
    if not (request.user.is_superuser or request.user.groups.filter(name='admin').exists()):
        return HttpResponseForbidden()
    b = get_object_or_404(Booking, pk=pk)
    days = b.days.order_by('day_no')
    return render(request,'admin_app/booking_detail.html',{'booking':b,'days':days})

# Instructor app
@login_required
def instructor_dashboard(request):
    if is_instructor_user(request.user):
        return redirect("instructor_bookings")
    return redirect("home")

# helper for human label + bootstrap class
_STATUS_MAP = {
    "scheduled":        ("Scheduled", "badge bg-info text-dark"),
    "in_progress":      ("In progress", "badge bg-dark"),
    "awaiting_closure": ("Awaiting instructor closure", "badge", "background-color:#6f42c1;color:#fff"),
    "completed":        ("Completed", "badge bg-success"),
    "cancelled":        ("Cancelled", "badge bg-warning text-dark"),
}

@login_required
def user_profile(request):

    from django.contrib.auth import update_session_auth_hash

    if request.user.is_superuser:
        messages.error(request, "Superuser cannot edit profile here.")
        return redirect("app_admin_dashboard")

    try:
        personnel = request.user.personnel
    except Personnel.DoesNotExist:
        messages.error(request, "No personnel profile linked to this account.")
        return redirect("no_roles")

    if request.method == "POST":
        uform = UserProfileForm(request.POST, instance=request.user)
        pform = PersonnelProfileForm(request.POST, instance=personnel)

        if uform.is_valid() and pform.is_valid():

            user = uform.save(commit=False)

            new_email = uform.cleaned_data["email"].strip().lower()
            user.username = new_email
            user.email = new_email
            user.save()

            # 🟢 FIX: sync Personnel.email
            personnel.email = new_email

            # 🟢 NEW: update dyslexia mode (checkbox present? → True)
            personnel.dyslexia_mode = bool(request.POST.get("dyslexia_mode"))
            personnel.pastel_background = request.POST.get("pastel_background", "none")

            pform.save()  # saves other personnel fields

            update_session_auth_hash(request, user)

            messages.success(request, "Your profile has been updated.")
            return redirect("user_profile")


    else:
        uform = UserProfileForm(instance=request.user)
        pform = PersonnelProfileForm(instance=personnel)

    return render(request, "profile.html", {
        "uform": uform,
        "pform": pform,
    })

# Public attendance
def public_attendance(request, booking_day_id):
    bday = get_object_or_404(BookingDay, pk=booking_day_id)
    if request.method=='POST':
        form = AttendanceForm(request.POST)
        if form.is_valid():
            att = form.save(commit=False)
            att.booking_day = bday
            att.save()
            return render(request,'public/attendance_success.html',{'booking_day':bday})
    else:
        form = AttendanceForm()
    return render(request,'public/attendance.html',{'form':form,'booking_day':bday})

# API
@login_required
def api_locations_by_business(request):
    bid = request.GET.get('business')
    data=[]
    if bid:
        for x in TrainingLocation.objects.filter(business_id=bid).order_by('name'):
            data.append({'id': str(x.id), 'name': f"{x.name} — {x.address_line or ''} {x.postcode or ''}"})
    return JsonResponse({'data':data})

# Switch role view
@login_required
def switch_role(request, role):
    role = role.lower()
    # only allow switching to roles the user actually has
    if role == "admin" and (request.user.is_superuser or request.user.groups.filter(name__iexact="admin").exists()):
        request.session["active_role"] = "admin"
        messages.success(request, "Switched to Admin.")
    elif role == "instructor" and request.user.groups.filter(name__iexact="instructor").exists():
        request.session["active_role"] = "instructor"
        messages.success(request, "Switched to Instructor.")
    else:
        messages.error(request, "You don't have that role.")
    return redirect("home")

# Instructor my bookings view
@login_required
def instructor_bookings(request):
    # Require the user to be linked to an Instructor record
    try:
        instructor = Personnel.objects.select_related('user').get(user=request.user)
    except Instructor.DoesNotExist:
        messages.error(request, "Your user account isn’t linked to an Instructor profile yet.")
        return redirect('instructor_dashboard')  # or 'home'

    bookings = (
        Booking.objects
        .select_related('business', 'training_location', 'course_type')
        .filter(instructor=instructor)
        .order_by('-course_date', 'start_time')
    )
    return render(request, 'instructor/bookings.html', {
        'bookings': bookings,
        'instructor': instructor,
    })

@login_required
def instructor_booking_detail(request, pk):
    """
    Detail page for a single booking, visible only to its assigned instructor.
    """
    iinst = get_instructor_for_user(request.user)
    if not inst:
        return HttpResponseForbidden("Not an instructor account.")

    booking = get_object_or_404(
        Booking.objects.select_related("business", "training_location", "course_type", "instructor"),
        pk=pk,
    )
    if booking.instructor_id != inst.id:
        return HttpResponseForbidden("This booking is not assigned to you.")

    # Convenient status label
    status_labels = dict(getattr(Booking, "STATUS_CHOICES", []))
    status_label = status_labels.get(booking.status, booking.status)

    # Include the course days
    days = booking.days.all().order_by("date")

    return render(request, "instructor/booking_detail.html", {
        "title": "Booking details",
        "booking": booking,
        "status_label": status_label,
        "days": days,
        "today": timezone.localdate(),
    })

def _user_is_instructor(user):
    return hasattr(user, "personnel") and user.personnel.is_active and user.personnel.can_login


@login_required
def post_login_router(request):
    user = request.user

    if user.is_superuser or user.is_staff:
        return redirect("app_admin_dashboard")

    if _user_is_instructor(user):
        return redirect("instructor_dashboard")

    return redirect("home")



    # Fallback
    return redirect("app_admin_dashboard")


def _resolve_course_type(value):
    if not value:
        return None
    # Try UUID first
    try:
        return CourseType.objects.filter(id=uuid.UUID(str(value))).first()
    except Exception:
        pass
    # Fall back to course code (case-insensitive)
    return CourseType.objects.filter(code__iexact=str(value)).first()

def _parse_date_flexible(s: str):
    if not s:
        return None
    d = parse_date(s)  # expects YYYY-MM-DD
    if d:
        return d
    # Try dd/mm/yyyy
    if "/" in s:
        try:
            dd, mm, yyyy = s.split("/")
            return parse_date(f"{yyyy}-{int(mm):02d}-{int(dd):02d}")
        except Exception:
            return None
    return None

def public_feedback_instructors_api(request):
    """
    Returns instructors who are delivering the given course on the given date.
    Query params:
      - course: course code (e.g. FAAW) or course UUID
      - date:   YYYY-MM-DD (required for filtering)
    If either parameter is missing, returns an empty list (no noisy fallback).
    """
    course_param = (request.GET.get("course") or "").strip()
    date_str     = (request.GET.get("date") or "").strip()

    # Resolve course by code OR UUID
    ct = None
    if course_param:
        ct = CourseType.objects.filter(code__iexact=course_param).first()
        if not ct:
            try:
                uuid.UUID(course_param)
                ct = CourseType.objects.filter(pk=course_param).first()
            except Exception:
                ct = None

    # Parse date flexibly: YYYY-MM-DD or dd/mm/yyyy
    the_date = None
    if date_str:
        the_date = parse_date(date_str)
        if not the_date and "/" in date_str:
            try:
                dd, mm, yyyy = date_str.split("/")
                the_date = parse_date(f"{int(yyyy):04d}-{int(mm):02d}-{int(dd):02d}")
            except Exception:
                the_date = None

    # Strict behaviour: only return matches when both pieces are present
    qs = Personnel.objects.none()
    if ct and the_date:
        qs = (
            Personnel.objects
            .filter(
                bookings__course_type=ct,
                bookings__days__date=the_date,
            )
            .distinct()
            .order_by("name")
        )

    return JsonResponse({
        "options": [{"id": str(i.id), "name": i.name} for i in qs]
    })


def _resolve_course_type(q):
    if not q:
        return None
    q = q.strip()
    ct = CourseType.objects.filter(code__iexact=q).first()
    if ct:
        return ct
    try:
        uuid.UUID(q)
        return CourseType.objects.filter(pk=q).first()
    except Exception:
        return None

def _parse_flexible_date(s):
    if not s:
        return None
    d = parse_date(s)
    if d:
        return d
    if "/" in s:
        try:
            dd, mm, yyyy = s.split("/")
            return parse_date(f"{int(yyyy):04d}-{int(mm):02d}-{int(dd):02d}")
        except Exception:
            return None
    return None

def public_feedback_form(request):
    """
    Public course feedback form. Supports optional prefill via:
      ?course=<code or UUID>&date=YYYY-MM-DD&instructor=<UUID>
    """
    course_q = request.GET.get("course") or ""
    prefilled_course = _resolve_course_type(course_q)

    # Initials: date (today unless provided) and instructor (optional)
    init = {
        "date": _parse_flexible_date(request.GET.get("date") or "") or timezone.localdate()
    }
    inst_q = request.GET.get("instructor") or ""
    if inst_q:
        try:
            init["instructor"] = Personnel.objects.get(pk=inst_q)
        except Instructor.DoesNotExist:
            pass

    form = FeedbackForm(request.POST or None, initial=init)

    # If course isn’t prefilled, user must choose it; otherwise set hidden initial
    form.fields["course_type"].required = not bool(prefilled_course)
    if prefilled_course:
        form.fields["course_type"].initial = prefilled_course.id

    if request.method == "POST" and form.is_valid():
        cd = form.cleaned_data
        course = prefilled_course or cd.get("course_type")

        # Create the feedback row
        FeedbackResponse.objects.create(
            course_type = course,
            date        = cd.get("date"),
            instructor  = cd.get("instructor"),

            # ratings
            overall_rating       = cd.get("overall_rating"),
            prior_knowledge      = cd.get("prior_knowledge"),
            post_knowledge       = cd.get("post_knowledge"),
            q_purpose_clear      = cd.get("q_purpose_clear"),
            q_personal_needs     = cd.get("q_personal_needs"),
            q_exercises_useful   = cd.get("q_exercises_useful"),
            q_structure          = cd.get("q_structure"),
            q_pace               = cd.get("q_pace"),
            q_content_clear      = cd.get("q_content_clear"),
            q_instructor_knowledge = cd.get("q_instructor_knowledge"),
            q_materials_quality  = cd.get("q_materials_quality"),
            q_books_quality      = cd.get("q_books_quality"),
            q_venue_suitable     = cd.get("q_venue_suitable"),
            q_benefit_at_work    = cd.get("q_benefit_at_work"),
            q_benefit_outside    = cd.get("q_benefit_outside"),

            # free text + contact
            comments       = cd.get("comments") or "",
            wants_callback = cd.get("wants_callback") or False,
            contact_name   = cd.get("contact_name") or "",
            contact_email  = cd.get("contact_email") or "",
            contact_phone  = cd.get("contact_phone") or "",
        )

        messages.success(request, "Thanks for your feedback!")
        return redirect("public_feedback_thanks")

    return render(
        request,
        "public/feedback_form.html",
        {"form": form, "course_type": prefilled_course},
    )


def public_feedback_pdf(request, pk):
    """Minimal single-response PDF so the route works."""
    fb = get_object_or_404(
        FeedbackResponse.objects.select_related("course_type", "instructor"),
        pk=pk
    )

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=portrait(A4))
    W, H = portrait(A4)
    left, top = 15*mm, H - 15*mm

    y = top
    c.setFont("Helvetica-Bold", 12)
    c.drawCentredString(W/2, y, f"{fb.course_type.code} — Course Feedback")
    y -= 8*mm
    c.setFont("Helvetica", 10)
    c.drawString(left, y, f"Date: {fb.date.strftime('%d %b %Y')}")
    y -= 5*mm
    c.drawString(left, y, f"Instructor: {fb.instructor.name if fb.instructor else '—'}")
    y -= 10*mm

    c.setFont("Helvetica", 9)
    c.drawString(left, y, f"Prior knowledge: {fb.prior_knowledge or '—'} / 5")
    y -= 5*mm
    c.drawString(left, y, f"Post knowledge:  {fb.post_knowledge  or '—'} / 5")
    y -= 8*mm
    c.drawString(left, y, f"Comments: {(fb.comments or '')[:300]}")
    y -= 12*mm

    c.setFont("Helvetica", 8)
    c.drawString(left, 10*mm, "Unicorn Fire & Safety Solutions, Unicorn House, 6 Salendine, Shrewsbury, SY1 3XJ: info@unicornsafety.co.uk: 01743 360211")

    c.save()
    buf.seek(0)
    return FileResponse(buf, as_attachment=True, filename=f"feedback_{fb.course_type.code}_{fb.date:%Y%m%d}.pdf")

def public_feedback_thanks(request):
    """Simple thank-you page after feedback submission."""
    return render(request, "public/feedback_thanks.html")

def no_roles_assigned(request):
    return render(request, "no_roles.html")


@login_required
def accident_report_list(request):
    reports = AccidentReport.objects.order_by("-date", "-time")
    return render(request, "instructor/accident_report_list.html", {"reports": reports})

@login_required
def accident_report_create(request):
    if request.method == "POST":
        form = AccidentReportForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, "Accident report submitted successfully.")
            return redirect("accident_report_list")
    else:
        form = AccidentReportForm()

    return render(request, "public/accident_report_form.html", {"form": form})

from django.utils import timezone

def accident_report_public(request):
    form = AccidentReportForm(request.POST or None)
    if request.method == "POST" and form.is_valid():
        form.save()
        messages.success(request, "Accident report submitted successfully.")
        return redirect("accident_report_thanks")  # or wherever your thank-you lives

    return render(
        request,
        "public/accident_report_form.html",
        {
            "form": form,
            "GOOGLE_MAPS_API_KEY": settings.GOOGLE_MAPS_API_KEY,  # <-- important
        },
    )

def accident_report_thanks(request):
    return render(request, "public/accident_report_thanks.html")

# --- Accident reports: DETAIL -------------------------------------------------

@login_required
def accident_report_detail(request, pk):
    # Use select_related / only as you like; keeping simple
    report = get_object_or_404(AccidentReport, pk=pk)
    return render(request, "instructor/accident_report_detail.html", {"report": report})

# --- Accident reports: EXPORT PPTX (selected IDs) -----------------------------

@login_required
def accident_report_export_pptx(request):
    """
    POST with ids[]=<uuid>… -> returns a styled .pptx
    One 4:3 slide per report with subtle first-aid watermark + corner L accent.
    """
    if request.method != "POST":
        return HttpResponseBadRequest("POST only")

    ids = request.POST.getlist("ids[]") or request.POST.getlist("ids")
    if not ids:
        return HttpResponseBadRequest("No IDs supplied")

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

    reports = list(
        AccidentReport.objects.filter(pk__in=ids).order_by("date", "time")
    )
    if not reports:
        return HttpResponseBadRequest("No matching reports")

    # ----- Theme colours
    BRAND = {
        "primary": RGBColor(220, 53, 69),    # red title bar / accents
        "ink":     RGBColor(33, 37, 41),     # near-black text
        "chipbg":  RGBColor(243, 244, 246),  # pill bg
        "chipfg":  RGBColor(33, 37, 41),
        "rule":    RGBColor(222, 226, 230),  # rules
        "panelbg": RGBColor(248, 249, 250),  # main card
        "panelbd": RGBColor(230, 232, 235),  # panel border
    }

    prs = Presentation()
    # Force 4:3 (10" x 7.5")
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---------- helpers

    def add_title_bar(slide, title_text):
        bar = slide.shapes.add_shape(
            SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10.0), Inches(0.9)
        )
        bar.fill.solid(); bar.fill.fore_color.rgb = BRAND["primary"]
        bar.line.fill.background()
        tf = bar.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28); p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        tf.margin_left = Inches(0.4)

    def meta_pill(slide, text):
        pill = slide.shapes.add_shape(
            SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.0), Inches(9.2), Inches(0.65)
        )
        pill.fill.solid(); pill.fill.fore_color.rgb = BRAND["chipbg"]
        pill.line.fill.background()
        tf = pill.text_frame; tf.clear(); tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16); p.font.color.rgb = BRAND["chipfg"]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        return pill

    def section(slide, left, top, title, body, width, height):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame; tf.clear(); tf.word_wrap = True
        tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)

        h = tf.paragraphs[0]
        h.text = title
        h.font.size = Pt(12); h.font.bold = True; h.font.color.rgb = BRAND["ink"]

        txt = (body or "—").strip()
        L = len(txt)
        base = 16 if L < 300 else (14 if L < 600 else 12)

        p = tf.add_paragraph(); p.text = txt
        p.font.size = Pt(base); p.font.color.rgb = BRAND["ink"]
        return box

    def vline(slide, x, y, h):
        ln = slide.shapes.add_shape(SHAPE.RECTANGLE, x, y, Inches(0.02), h)
        ln.fill.solid(); ln.fill.fore_color.rgb = BRAND["rule"]
        ln.line.fill.background()

    def watermark_plus(slide):
        """
        Very faint first-aid '+' watermark in the canvas centre.
        """
        cx, cy = Inches(5.0), Inches(4.1)
        w, t = Inches(2.8), Inches(0.45)   # arm length / thickness

        # horizontal
        hbar = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, cx - w/2, cy - t/2, w, t)
        # vertical
        vbar = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, cx - t/2, cy - w/2, t, w)
        for shp in (hbar, vbar):
            shp.fill.solid()
            shp.fill.fore_color.rgb = BRAND["primary"]
            shp.fill.fore_color.brightness = 0.7   # push toward white
            shp.line.fill.background()

    def corner_L_accent(slide):
        """
        Soft brand L-shape in the top-left of the inner panel.
        """
        # short horizontal
        h = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(0.25), Inches(0.95), Inches(0.9), Inches(0.07))
        # long vertical
        v = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(0.25), Inches(0.95), Inches(0.07), Inches(1.3))
        for shp in (h, v):
            shp.fill.solid()
            shp.fill.fore_color.rgb = BRAND["primary"]
            shp.fill.fore_color.brightness = 0.5   # softer than title
            shp.line.fill.background()

    # ---------- build slides

    for ar in reports:
        slide = prs.slides.add_slide(blank)

        # Background card
        panel = slide.shapes.add_shape(
            SHAPE.RECTANGLE, Inches(0.2), Inches(0.9), Inches(9.6), Inches(6.3)
        )
        panel.fill.solid(); panel.fill.fore_color.rgb = BRAND["panelbg"]
        panel.line.color.rgb = BRAND["panelbd"]

        # Motifs first (behind text visually)
        watermark_plus(slide)
        corner_L_accent(slide)

        # Header
        add_title_bar(slide, "Accident / Incident Report")

        # Meta pill
        meta_text = " • ".join([
            f"Date: {ar.date:%d %b %Y}" if ar.date else "Date: —",
            f"Time: {ar.time.strftime('%H:%M')}" if ar.time else "Time: —",
            f"Location: {ar.location or '—'}",
        ])
        meta_pill(slide, meta_text)

        # Columns
        left_x  = Inches(0.45)
        right_x = Inches(5.1)
        top_y   = Inches(1.9)

        # Left: people + address
        section(
            slide, left_x, top_y,
            "People involved",
            f"Injured: {ar.injured_name or '—'}\n"
            f"First aider: {ar.first_aider_name or '—'}\n"
            f"Reporter: {ar.reporter_name or '—'}",
            width=Inches(4.4), height=Inches(1.45),
        )
        section(
            slide, left_x, top_y + Inches(1.6),
            "Injured address",
            ar.injured_address or "—",
            width=Inches(4.4), height=Inches(1.1),
        )

        # Divider
        vline(slide, Inches(4.95), Inches(1.85), Inches(4.45))

        # Right: narratives
        ry = top_y
        section(slide, right_x, ry,                   "What happened",               ar.what_happened,             width=Inches(4.35), height=Inches(1.35))
        section(slide, right_x, ry + Inches(1.5),     "Injuries sustained",         ar.injuries_sustained,        width=Inches(4.35), height=Inches(1.05))
        section(slide, right_x, ry + Inches(2.7),     "Actions carried out",        ar.actions_carried_out,       width=Inches(4.35), height=Inches(1.05))
        section(slide, right_x, ry + Inches(3.9),     "Actions to prevent recurrence", ar.actions_prevent_recurrence, width=Inches(4.35), height=Inches(1.05))

        # Footer tag
        footer = slide.shapes.add_textbox(Inches(0.3), Inches(7.25), Inches(9.4), Inches(0.3))
        ft = footer.text_frame; ft.clear()
        fp = ft.paragraphs[0]
        fp.text = "Confidential — Internal Safety Record"
        fp.font.size = Pt(10); fp.font.color.rgb = RGBColor(160, 165, 170)

    # ---------- return file
    from io import BytesIO
    buf = BytesIO(); prs.save(buf); buf.seek(0)
    resp = HttpResponse(
        buf.read(),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    resp["Content-Disposition"] = 'attachment; filename="Accident reports.pptx"'
    return resp

@require_POST
@login_required
def accident_report_delete(request):
    """
    Bulk-delete accident reports by ids[]=...
    """
    ids = request.POST.getlist("ids[]") or request.POST.getlist("ids")
    if not ids:
        messages.warning(request, "No reports selected.")
        return redirect("accident_report_list")

    qs = AccidentReport.objects.filter(pk__in=ids)
    count = qs.count()
    if count == 0:
        messages.warning(request, "No matching reports found.")
        return redirect("accident_report_list")

    qs.delete()
    messages.success(request, f"Deleted {count} report{'s' if count != 1 else ''}.")
    return redirect("accident_report_list")

@login_required
def accident_report_poll(request):
    """
    Returns just the table rows (HTML) for the accident report list.
    Frontend replaces <tbody> with this HTML, preserving selections.
    """
    reports = AccidentReport.objects.order_by("-date", "-time", "-id")
    rows_html = render_to_string(
        "instructor/_accident_report_rows.html",
        {"reports": reports},
        request=request,
    )
    return JsonResponse({"html": rows_html})

@require_POST
def accident_report_anonymise(request):
    ids = request.POST.getlist("ids")  # matches name="ids" in the list form rows
    if not ids:
        messages.warning(request, "Select at least one report to anonymise.")
        return redirect("accident_report_list")

    qs = AccidentReport.objects.filter(id__in=ids, anonymized_at__isnull=True)
    # Use update for speed, then set anonymized_at in a second pass
    now = timezone.now()
    updated = qs.update(injured_name="", injured_address="", anonymized_at=now)

    messages.success(request, f"Anonymised {updated} report{'s' if updated != 1 else ''}.")
    return redirect("accident_report_list")